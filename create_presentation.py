#!/usr/bin/env python3
"""
Generiert eine professionelle PowerPoint-Präsentation für das Projekt "most wanTED".
FH Campus 02 – Integrationstechnologien für eGovernment
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── Farbpalette (dunkles Theme) ────────────────────────────────────────
BG_DARK       = RGBColor(0x0F, 0x17, 0x2A)   # Sehr dunkles Blau
BG_CARD       = RGBColor(0x1A, 0x25, 0x3C)   # Karten-Hintergrund
ACCENT_BLUE   = RGBColor(0x3B, 0x82, 0xF6)   # Primärfarbe Blau
ACCENT_CYAN   = RGBColor(0x06, 0xB6, 0xD4)   # Sekundärfarbe Cyan
ACCENT_GREEN  = RGBColor(0x10, 0xB9, 0x81)   # Grün für Erfolg
ACCENT_ORANGE = RGBColor(0xF5, 0x9E, 0x0B)   # Orange für Warnung
ACCENT_PURPLE = RGBColor(0x8B, 0x5C, 0xF6)   # Lila für Highlights
ACCENT_RED    = RGBColor(0xEF, 0x44, 0x44)   # Rot für Fehler
TEXT_WHITE    = RGBColor(0xF1, 0xF5, 0xF9)   # Heller Text
TEXT_GRAY     = RGBColor(0x94, 0xA3, 0xB8)   # Grauer Subtext
TEXT_LIGHT    = RGBColor(0xCB, 0xD5, 0xE1)   # Leicht grauer Text
DIVIDER       = RGBColor(0x33, 0x41, 0x55)   # Trennlinie


def set_slide_bg(slide, color=BG_DARK):
    """Setzt den Hintergrund einer Folie."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(0)):
    """Fügt ein Rechteck hinzu."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_accent_bar(slide, left, top, width, height, color=ACCENT_BLUE):
    """Schmaler Akzentbalken."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=14,
                 color=TEXT_WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    """Fügt eine Textbox hinzu."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=13, color=TEXT_LIGHT, spacing=Pt(6)):
    """Fügt eine Bullet-Liste hinzu."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Segoe UI"
        p.space_after = spacing
        p.level = 0
    return txBox


def add_table_slide(slide, left, top, width, rows_data, col_widths, header_color=ACCENT_BLUE):
    """Fügt eine formatierte Tabelle hinzu."""
    cols = len(rows_data[0])
    rows = len(rows_data)
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, Inches(0.4 * rows))
    table = table_shape.table

    for i, cw in enumerate(col_widths):
        table.columns[i].width = cw

    for row_idx, row_data in enumerate(rows_data):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(cell_text)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.font.name = "Segoe UI"
                if row_idx == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = TEXT_WHITE
                else:
                    paragraph.font.color.rgb = TEXT_LIGHT
                paragraph.alignment = PP_ALIGN.LEFT

            if row_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            elif row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0x1E, 0x29, 0x3B)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = BG_CARD

    return table_shape


def add_slide_number(slide, num, total):
    """Fügt Foliennummer hinzu."""
    add_text_box(slide, Inches(8.8), Inches(7.0), Inches(1.0), Inches(0.3),
                 f"{num}/{total}", font_size=9, color=TEXT_GRAY, alignment=PP_ALIGN.RIGHT)


def add_section_header(slide, title, subtitle="", accent_color=ACCENT_BLUE):
    """Standard-Folienheader mit Akzentbalken."""
    add_accent_bar(slide, Inches(0.5), Inches(0.4), Inches(0.08), Inches(0.55), accent_color)
    add_text_box(slide, Inches(0.75), Inches(0.35), Inches(8.5), Inches(0.5),
                 title, font_size=26, color=TEXT_WHITE, bold=True)
    if subtitle:
        add_text_box(slide, Inches(0.75), Inches(0.85), Inches(8.5), Inches(0.3),
                     subtitle, font_size=13, color=TEXT_GRAY)


def add_footer_line(slide):
    """Dezente Fußleiste."""
    add_accent_bar(slide, Inches(0.5), Inches(6.85), Inches(9.0), Pt(1), DIVIDER)
    add_text_box(slide, Inches(0.5), Inches(6.9), Inches(4.0), Inches(0.3),
                 "most wanTED — FH Campus 02", font_size=8, color=TEXT_GRAY)


# ═══════════════════════════════════════════════════════════════════════
#  PRÄSENTATION ERSTELLEN
# ═══════════════════════════════════════════════════════════════════════

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

TOTAL_SLIDES = 20

# ──────────────────────────────────────────────────────────────────────
# FOLIE 1: Titelfolie
# ──────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
set_slide_bg(slide)

# Gradient-ähnlicher Akzent oben
add_accent_bar(slide, Inches(0), Inches(0), Inches(10), Pt(4), ACCENT_BLUE)

# Titel
add_text_box(slide, Inches(0.8), Inches(1.8), Inches(8.4), Inches(1.0),
             "most wanTED", font_size=52, color=TEXT_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Untertitel
add_text_box(slide, Inches(0.8), Inches(2.9), Inches(8.4), Inches(0.6),
             "Microservice-Architektur für die Trading Platform", font_size=22,
             color=ACCENT_CYAN, alignment=PP_ALIGN.CENTER)

# Trennlinie
add_accent_bar(slide, Inches(3.5), Inches(3.7), Inches(3.0), Pt(2), ACCENT_BLUE)

# Team & Kontext
add_text_box(slide, Inches(0.8), Inches(4.1), Inches(8.4), Inches(0.4),
             "Hans Erik Krenn  ·  Patrick Grüner  ·  Kevin Ulm", font_size=16,
             color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(0.8), Inches(4.6), Inches(8.4), Inches(0.4),
             "FH Campus 02 — Integrationstechnologien für eGovernment", font_size=13,
             color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(0.8), Inches(5.1), Inches(8.4), Inches(0.4),
             "SS 2026", font_size=13,
             color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

# Tech-Badge unten
for i, tech in enumerate([".NET 8", "C#", "Consul", "gRPC", "OData", "Polly"]):
    x = Inches(1.6 + i * 1.15)
    card = add_shape(slide, x, Inches(5.8), Inches(1.0), Inches(0.35), BG_CARD, ACCENT_BLUE, Pt(1))
    add_text_box(slide, x, Inches(5.82), Inches(1.0), Inches(0.35),
                 tech, font_size=10, color=ACCENT_CYAN, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 1, TOTAL_SLIDES)

# ──────────────────────────────────────────────────────────────────────
# FOLIE 2: Agenda
# ──────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_header(slide, "Agenda", "Überblick über die Präsentation")

agenda_items = [
    ("01", "Projektziel & Kontext", ACCENT_BLUE),
    ("02", "Architektur-Übersicht", ACCENT_CYAN),
    ("03", "Service-Landschaft", ACCENT_GREEN),
    ("04", "Kommunikationsprotokolle", ACCENT_PURPLE),
    ("05", "Domain Driven Design", ACCENT_ORANGE),
    ("06", "Patterns & Konzepte", ACCENT_BLUE),
    ("07", "Live-Demo Szenarien", ACCENT_CYAN),
    ("08", "Vertiefungsthemen", ACCENT_GREEN),
    ("09", "Fazit & Ausblick", ACCENT_PURPLE),
]

for i, (num, title, color) in enumerate(agenda_items):
    row = i // 3
    col = i % 3
    x = Inches(0.6 + col * 3.1)
    y = Inches(1.6 + row * 1.7)
    card = add_shape(slide, x, y, Inches(2.8), Inches(1.3), BG_CARD, color, Pt(1))
    add_text_box(slide, x + Inches(0.15), y + Inches(0.15), Inches(2.5), Inches(0.4),
                 num, font_size=28, color=color, bold=True)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.7), Inches(2.5), Inches(0.5),
                 title, font_size=14, color=TEXT_WHITE)

add_footer_line(slide)
add_slide_number(slide, 2, TOTAL_SLIDES)

# ──────────────────────────────────────────────────────────────────────
# FOLIE 3: Projektziel & Kontext
# ──────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_header(slide, "Projektziel & Kontext", 'Was ist "most wanTED"?')

# Haupttext
add_text_box(slide, Inches(0.75), Inches(1.5), Inches(8.5), Inches(1.2),
             "most wanTED erweitert die bestehende Handelsplattform SolTradingPlatform um "
             "eine umfragebasierte Entscheidungsunterstützung. Beteiligte — Kunden, Lieferanten, "
             "Kreditkartenunternehmen und Produzenten — werden bei Entscheidungen zu Produktauswahl, "
             "Bezahlvarianten und Liefervarianten unterstützt.",
             font_size=14, color=TEXT_LIGHT)

# Info-Karten
cards = [
    ("🎓 Lehrveranstaltung", "Integrationstechnologien\nfür eGovernment\nFH Campus 02, Graz", ACCENT_BLUE),
    ("👥 Team", "Hans Erik Krenn\nPatrick Grüner\nKevin Ulm", ACCENT_CYAN),
    ("⚙️ Tech-Stack", ".NET 8 / C#\nASP.NET Core\nMicroservices", ACCENT_GREEN),
    ("📊 Bewertung", "100 Punkte regulär\n+ 5 Bonus\n= 105 möglich", ACCENT_ORANGE),
]

for i, (title, desc, color) in enumerate(cards):
    x = Inches(0.55 + i * 2.3)
    y = Inches(3.2)
    card = add_shape(slide, x, y, Inches(2.1), Inches(2.8), BG_CARD, color, Pt(1))
    add_accent_bar(slide, x, y, Inches(2.1), Pt(3), color)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.2), Inches(1.8), Inches(0.4),
                 title, font_size=13, color=color, bold=True)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.7), Inches(1.8), Inches(2.0),
                 desc, font_size=12, color=TEXT_LIGHT)

add_footer_line(slide)
add_slide_number(slide, 3, TOTAL_SLIDES)

# ──────────────────────────────────────────────────────────────────────
# FOLIE 4: Architektur-Übersicht (Makro)
# ──────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_header(slide, "Makro-Architektur", "Service-Landschaft & Kommunikationsbeziehungen")

# Client/Browser Box
add_shape(slide, Inches(3.8), Inches(1.4), Inches(2.4), Inches(0.6), BG_CARD, TEXT_GRAY, Pt(1))
add_text_box(slide, Inches(3.8), Inches(1.44), Inches(2.4), Inches(0.5),
             "🌐 Browser / Client", font_size=13, color=TEXT_WHITE, alignment=PP_ALIGN.CENTER)

# Pfeil
add_text_box(slide, Inches(4.7), Inches(2.0), Inches(0.6), Inches(0.3),
             "▼", font_size=16, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

# API Gateway (MeiShop)
gw = add_shape(slide, Inches(2.8), Inches(2.35), Inches(4.4), Inches(0.7), RGBColor(0x1E, 0x3A, 0x5F), ACCENT_BLUE, Pt(2))
add_text_box(slide, Inches(2.8), Inches(2.4), Inches(4.4), Inches(0.6),
             "MeiShop — API Gateway (Port 7024)", font_size=15, color=ACCENT_CYAN, bold=True, alignment=PP_ALIGN.CENTER)

# Backend Services
services = [
    ("ProductService\n:7200", ACCENT_GREEN),
    ("FtpCatalog\n:7300", ACCENT_GREEN),
    ("CreditCard\n:7231-33", ACCENT_ORANGE),
    ("PaymentSvc\n:7400", ACCENT_ORANGE),
    ("ODataSvc\n:7500", ACCENT_PURPLE),
    ("SagaSvc\n:7700", ACCENT_RED),
]

for i, (name, color) in enumerate(services):
    x = Inches(0.3 + i * 1.55)
    y = Inches(3.5)
    card = add_shape(slide, x, y, Inches(1.4), Inches(1.1), BG_CARD, color, Pt(1))
    add_text_box(slide, x, y + Inches(0.1), Inches(1.4), Inches(0.9),
                 name, font_size=10, color=TEXT_WHITE, alignment=PP_ALIGN.CENTER)

# Infrastruktur-Services
infra = [
    ("Consul — Service Discovery (:8500)", ACCENT_BLUE),
    ("LoggingService — gRPC (:5500)", ACCENT_CYAN),
    ("WebhookSubscriber (:7600)", ACCENT_GREEN),
]

for i, (name, color) in enumerate(infra):
    x = Inches(0.5 + i * 3.1)
    y = Inches(5.2)
    card = add_shape(slide, x, y, Inches(2.9), Inches(0.55), BG_CARD, color, Pt(1))
    add_text_box(slide, x, y + Inches(0.05), Inches(2.9), Inches(0.45),
                 name, font_size=10, color=color, alignment=PP_ALIGN.CENTER)

# Legende
add_text_box(slide, Inches(0.5), Inches(6.1), Inches(9.0), Inches(0.4),
             "REST/JSON  ·  gRPC/Protobuf  ·  OData v4  ·  Webhooks  ·  Content Negotiation (JSON/XML/CSV)",
             font_size=10, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

add_footer_line(slide)
add_slide_number(slide, 4, TOTAL_SLIDES)

# ──────────────────────────────────────────────────────────────────────
# FOLIE 5: Service-Übersicht (Tabelle)
# ──────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_header(slide, "Service-Übersicht", "Alle Microservices im Überblick")

table_data = [
    ["Service", "Port", "Protokoll", "Beschreibung"],
    ["MeiShop", "7024", "REST (MVC)", "Frontend / API Gateway"],
    ["ProductService", "7200", "REST (JSON)", "Produktverwaltung (CRUD)"],
    ["FtpProductCatalog", "7300", "REST (JSON)", "Produkte aus Textdatei (FTP)"],
    ["IEGEasyCreditCard", "7231–33", "REST (JSON)", "Kreditkarte (3 Instanzen)"],
    ["PaymentService", "7400", "REST (JSON/XML/CSV)", "Payment + Content Negotiation"],
    ["LoggingService", "5500", "gRPC (Protobuf)", "Zentrales Logging"],
    ["ProductODataService", "7500", "OData v4", "Standardisierte Abfragen"],
    ["WebhookSubscriber", "7600", "REST (JSON)", "Webhook-Empfänger"],
    ["OrderSagaService", "7700", "REST (JSON)", "SAGA-Orchestrator"],
    ["Consul", "8500", "HTTP API", "Service Discovery"],
]

add_table_slide(slide, Inches(0.5), Inches(1.5), Inches(9.0), table_data,
                [Inches(2.2), Inches(1.0), Inches(2.2), Inches(3.6)])

add_footer_line(slide)
add_slide_number(slide, 5, TOTAL_SLIDES)

# ──────────────────────────────────────────────────────────────────────
# FOLIE 6: Kommunikationsprotokolle
# ──────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_header(slide, "Kommunikationsprotokolle", "Wie kommunizieren die Services?")

protocols = [
    ("REST / HTTP", "JSON", "Primäre Service-zu-Service\nKommunikation. Standard\nASP.NET Core Controller.", ACCENT_BLUE),
    ("gRPC", "Protobuf", "Binäres Protokoll für\neffizientes Logging.\nHoher Durchsatz.", ACCENT_CYAN),
    ("OData v4", "JSON", "SQL-ähnliche Abfragen:\n$filter, $orderby, $top\n$skip, $select, $count", ACCENT_GREEN),
    ("Webhooks", "JSON", "Event-basierte Push-\nBenachrichtigungen bei\nProduktänderungen.", ACCENT_ORANGE),
    ("Content\nNegotiation", "JSON\nXML\nCSV", "Format-Auswahl über\nHTTP Accept/Content-Type\nHeader.", ACCENT_PURPLE),
]

for i, (title, fmt, desc, color) in enumerate(protocols):
    x = Inches(0.3 + i * 1.9)
    y = Inches(1.6)
    h = Inches(4.8)
    card = add_shape(slide, x, y, Inches(1.75), h, BG_CARD, color, Pt(1))
    add_accent_bar(slide, x, y, Inches(1.75), Pt(3), color)
    add_text_box(slide, x + Inches(0.1), y + Inches(0.2), Inches(1.55), Inches(0.7),
                 title, font_size=14, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    # Format badge
    badge = add_shape(slide, x + Inches(0.2), y + Inches(1.0), Inches(1.35), Inches(0.7), 
                      RGBColor(0x0F, 0x17, 0x2A), color, Pt(1))
    add_text_box(slide, x + Inches(0.2), y + Inches(1.05), Inches(1.35), Inches(0.65),
                 fmt, font_size=11, color=color, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), y + Inches(2.0), Inches(1.55), Inches(2.5),
                 desc, font_size=11, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

add_footer_line(slide)
add_slide_number(slide, 6, TOTAL_SLIDES)

# ──────────────────────────────────────────────────────────────────────
# FOLIE 7: Mikro-Architektur — PaymentService
# ──────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_header(slide, "Mikro-Architektur: PaymentService", "Interner Aufbau eines Microservices")

# Schichten
layers = [
    ("Controllers/", "PaymentsController.cs\nWebhookSubscriptionsController.cs", ACCENT_BLUE, 
     "REST-API Endpunkte\nContent Negotiation"),
    ("Formatters/", "CsvInputFormatter.cs\nCsvOutputFormatter.cs", ACCENT_CYAN,
     "Custom Format-Handler\nfür CSV-Daten"),
    ("Services/", "PaymentRepository.cs\nWebhookService.cs", ACCENT_GREEN,
     "Geschäftslogik &\nDatenzugriff"),
    ("Models/", "Payment.cs\nWebhookSubscription.cs", ACCENT_ORANGE,
     "Datenmodelle"),
]

for i, (folder, files, color, desc) in enumerate(layers):
    y = Inches(1.5 + i * 1.3)
    # Folder name
    card = add_shape(slide, Inches(0.5), y, Inches(1.8), Inches(1.0), BG_CARD, color, Pt(1))
    add_accent_bar(slide, Inches(0.5), y, Pt(4), Inches(1.0), color)
    add_text_box(slide, Inches(0.65), y + Inches(0.1), Inches(1.6), Inches(0.35),
                 folder, font_size=14, color=color, bold=True)
    add_text_box(slide, Inches(0.65), y + Inches(0.45), Inches(1.6), Inches(0.5),
                 files, font_size=9, color=TEXT_GRAY)
    # Arrow
    add_text_box(slide, Inches(2.4), y + Inches(0.2), Inches(0.4), Inches(0.5),
                 "→", font_size=20, color=color, alignment=PP_ALIGN.CENTER)
    # Description
    add_shape(slide, Inches(2.9), y, Inches(6.5), Inches(1.0), BG_CARD)
    add_text_box(slide, Inches(3.05), y + Inches(0.15), Inches(6.2), Inches(0.7),
                 desc, font_size=13, color=TEXT_LIGHT)

add_footer_line(slide)
add_slide_number(slide, 7, TOTAL_SLIDES)

# ──────────────────────────────────────────────────────────────────────
# FOLIE 8: Domain Driven Design
# ──────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_header(slide, "Domain Driven Design", "Bounded Contexts & Beziehungsmuster")

contexts = [
    ("Gateway Context", "MeiShop", ACCENT_BLUE, "UI, Routing,\nAggregation"),
    ("Catalog Context", "ProductService\nFtpCatalog\nODataService", ACCENT_GREEN, "Produktverwaltung\n& -abfrage"),
    ("Payment Context", "PaymentService\nCreditCardSvc", ACCENT_ORANGE, "Zahlungs-\nverarbeitung"),
    ("Order Context", "OrderSaga\nService", ACCENT_RED, "Bestell-\norchestration"),
    ("Infrastructure", "LoggingService\nConsul\nWebhookSub", ACCENT_PURPLE, "Querschnitts-\nbelange"),
]

for i, (ctx, svcs, color, desc) in enumerate(contexts):
    x = Inches(0.3 + i * 1.9)
    y = Inches(1.6)
    card = add_shape(slide, x, y, Inches(1.75), Inches(2.8), BG_CARD, color, Pt(2))
    add_accent_bar(slide, x, y, Inches(1.75), Pt(3), color)
    add_text_box(slide, x + Inches(0.05), y + Inches(0.2), Inches(1.65), Inches(0.5),
                 ctx, font_size=12, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.05), y + Inches(0.7), Inches(1.65), Inches(1.0),
                 svcs, font_size=10, color=TEXT_WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.05), y + Inches(1.8), Inches(1.65), Inches(0.8),
                 desc, font_size=10, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

# DDD-Patterns Tabelle
ddd_data = [
    ["DDD-Konzept", "Einsatz"],
    ["Open Host Service", "ProductService, PaymentService — stabile öffentliche APIs"],
    ["Customer/Supplier", "MeiShop konsumiert Backend-Services"],
    ["Anticorruption Layer", "OrderSagaService ↔ ProductService, PaymentService"],
    ["Conformist", "WebhookSubscriber übernimmt Format 1:1"],
    ["Shared Kernel", "Bewusst NICHT eingesetzt — maximale Unabhängigkeit"],
]

add_table_slide(slide, Inches(0.5), Inches(4.7), Inches(9.0), ddd_data,
                [Inches(2.8), Inches(6.2)], ACCENT_PURPLE)

add_footer_line(slide)
add_slide_number(slide, 8, TOTAL_SLIDES)

# ──────────────────────────────────────────────────────────────────────
# FOLIE 9: Implementierte Patterns (Übersicht)
# ──────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_header(slide, "Implementierte Patterns", "Architekturelle Entscheidungen & Patterns")

patterns_data = [
    ["Pattern", "Service", "Beschreibung"],
    ["API Gateway", "MeiShop", "Zentraler Einstiegspunkt, routet zu Backend-Services"],
    ["Service Discovery", "Alle → Consul", "Dynamische Service-Lokalisierung via Consul HTTP API"],
    ["Round-Robin LB", "MeiShop → CreditCard", "Verteilt Anfragen gleichmäßig auf 3 Instanzen"],
    ["Retry / Backoff", "MeiShop (Polly)", "Automatisches Wiederholen bei transienten Fehlern"],
    ["Content Negotiation", "PaymentService", "JSON, XML, CSV via Accept/Content-Type"],
    ["Webhook Pub/Sub", "ProductSvc → WebhookSub", "Event-Benachrichtigung bei Produktänderungen"],
    ["SAGA Orchestration", "OrderSagaService", "Verteilte Transaktion mit Compensating Transactions"],
    ["Decentralized Data", "Alle Services", "Jeder Service hat eigenen In-Memory-Datastore"],
    ["Design for Failure", "CreditCard + MeiShop", "Simulierte Fehler + Retry + Fallback"],
    ["Multi-Instance", "IEGEasyCreditCard", "3 Instanzen auf verschiedenen Ports"],
]

add_table_slide(slide, Inches(0.3), Inches(1.4), Inches(9.4), patterns_data,
                [Inches(2.3), Inches(2.5), Inches(4.6)])

add_footer_line(slide)
add_slide_number(slide, 9, TOTAL_SLIDES)

# ──────────────────────────────────────────────────────────────────────
# FOLIE 10: Skalierung & Resilience
# ──────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_header(slide, "Skalierung & Resilience", "Design for Failure — IEGEasyCreditCardService")

# Flow diagram
steps = [
    ("MeiShop\nRequest", ACCENT_BLUE),
    ("Round Robin\nLoad Balancer", ACCENT_CYAN),
    ("Instanz 7231\n❌ Fehler", ACCENT_RED),
    ("Polly Retry\n(Exp. Backoff)", ACCENT_ORANGE),
    ("Instanz 7232\n✅ Erfolg", ACCENT_GREEN),
]

for i, (label, color) in enumerate(steps):
    x = Inches(0.3 + i * 1.95)
    y = Inches(1.6)
    card = add_shape(slide, x, y, Inches(1.7), Inches(1.2), BG_CARD, color, Pt(2))
    add_text_box(slide, x, y + Inches(0.15), Inches(1.7), Inches(0.9),
                 label, font_size=13, color=TEXT_WHITE, alignment=PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        add_text_box(slide, x + Inches(1.7), y + Inches(0.3), Inches(0.25), Inches(0.4),
                     "→", font_size=18, color=color)

# Details
add_text_box(slide, Inches(0.5), Inches(3.2), Inches(9.0), Inches(0.3),
             "Konfiguration", font_size=16, color=ACCENT_CYAN, bold=True)

config_items = [
    "• 3 Instanzen: Ports 7231, 7232, 7233 — Consul-Registration mit eindeutiger ID",
    "• Polly Retry: 4 Versuche pro Instanz mit exponentiellem Backoff (2s → 4s → 8s → 16s)",
    "• Failover: Nach 4 Fehlern → nächste Instanz im Round-Robin",
    "• gRPC Logging: Jeder Fehler wird an LoggingService gesendet (fire-and-forget)",
    "• Fehlersimulation: 20% zufällige Fehlerrate im CreditCardService für Demo-Zwecke",
]

add_bullet_list(slide, Inches(0.5), Inches(3.6), Inches(9.0), Inches(3.0), config_items,
                font_size=13, color=TEXT_LIGHT, spacing=Pt(8))

add_footer_line(slide)
add_slide_number(slide, 10, TOTAL_SLIDES)

# ──────────────────────────────────────────────────────────────────────
# FOLIE 11: SAGA Pattern
# ──────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_header(slide, "SAGA Pattern — OrderSagaService", "Verteilte Transaktionen mit Compensating Transactions")

# SAGA Steps
saga_steps = [
    ("1. Reserve\nProducts", "ProductService", "Kompensation:\nUnreserve", ACCENT_GREEN),
    ("2. Process\nPayment", "PaymentService", "Kompensation:\nRefund", ACCENT_ORANGE),
    ("3. Confirm\nOrder", "OrderSagaService", "Kompensation:\nCancel", ACCENT_BLUE),
]

for i, (step, svc, comp, color) in enumerate(saga_steps):
    x = Inches(0.5 + i * 3.2)
    y = Inches(1.6)
    # Main step
    card = add_shape(slide, x, y, Inches(2.8), Inches(1.3), BG_CARD, color, Pt(2))
    add_text_box(slide, x + Inches(0.1), y + Inches(0.1), Inches(2.6), Inches(0.7),
                 step, font_size=16, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), y + Inches(0.85), Inches(2.6), Inches(0.35),
                 f"→ {svc}", font_size=11, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)
    # Compensation
    comp_card = add_shape(slide, x + Inches(0.3), y + Inches(1.6), Inches(2.2), Inches(0.8), 
                          RGBColor(0x2D, 0x15, 0x15), ACCENT_RED, Pt(1))
    add_text_box(slide, x + Inches(0.3), y + Inches(1.65), Inches(2.2), Inches(0.7),
                 comp, font_size=11, color=ACCENT_RED, alignment=PP_ALIGN.CENTER)
    if i < 2:
        add_text_box(slide, x + Inches(2.8), y + Inches(0.35), Inches(0.4), Inches(0.4),
                     "→", font_size=22, color=color)

# State Machine
add_text_box(slide, Inches(0.5), Inches(4.1), Inches(9.0), Inches(0.3),
             "State Machine", font_size=16, color=ACCENT_CYAN, bold=True)

states = ["Pending", "ProductsReserved", "PaymentProcessed", "Confirmed"]
for i, state in enumerate(states):
    x = Inches(0.5 + i * 2.35)
    color = [TEXT_GRAY, ACCENT_GREEN, ACCENT_ORANGE, ACCENT_BLUE][i]
    card = add_shape(slide, x, Inches(4.6), Inches(2.0), Inches(0.5), BG_CARD, color, Pt(1))
    add_text_box(slide, x, Inches(4.63), Inches(2.0), Inches(0.45),
                 state, font_size=11, color=color, alignment=PP_ALIGN.CENTER)
    if i < 3:
        add_text_box(slide, x + Inches(2.0), Inches(4.6), Inches(0.35), Inches(0.5),
                     "→", font_size=14, color=TEXT_GRAY)

# Error path
error_states = ["Failed", "Compensated"]
for i, state in enumerate(error_states):
    x = Inches(3.5 + i * 2.35)
    card = add_shape(slide, x, Inches(5.4), Inches(2.0), Inches(0.5), BG_CARD, ACCENT_RED, Pt(1))
    add_text_box(slide, x, Inches(5.43), Inches(2.0), Inches(0.45),
                 state, font_size=11, color=ACCENT_RED, alignment=PP_ALIGN.CENTER)
    if i < 1:
        add_text_box(slide, x + Inches(2.0), Inches(5.4), Inches(0.35), Inches(0.5),
                     "→", font_size=14, color=ACCENT_RED)

add_footer_line(slide)
add_slide_number(slide, 11, TOTAL_SLIDES)

# ──────────────────────────────────────────────────────────────────────
# FOLIE 12: Content Negotiation
# ──────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_header(slide, "Content Negotiation", "PaymentService — JSON, XML & CSV")

# 3 Format-Beispiele
formats = [
    ("JSON", "application/json",
     '{\n  "Id": 1,\n  "Amount": 49.99,\n  "Currency": "EUR",\n  "Description": "Order",\n  "PaymentMethod": "Card"\n}',
     ACCENT_BLUE),
    ("XML", "application/xml",
     '<Payment>\n  <Id>1</Id>\n  <Amount>49.99</Amount>\n  <Currency>EUR</Currency>\n  <Description>Order</Description>\n  <PaymentMethod>Card</PaymentMethod>\n</Payment>',
     ACCENT_GREEN),
    ("CSV", "text/csv",
     'Id,Amount,Currency,\nDescription,PaymentMethod\n1,49.99,EUR,\nOrder,Card',
     ACCENT_ORANGE),
]

for i, (fmt, mime, example, color) in enumerate(formats):
    x = Inches(0.4 + i * 3.15)
    y = Inches(1.5)
    card = add_shape(slide, x, y, Inches(2.9), Inches(4.8), BG_CARD, color, Pt(1))
    add_accent_bar(slide, x, y, Inches(2.9), Pt(3), color)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.2), Inches(2.6), Inches(0.4),
                 fmt, font_size=22, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    # MIME badge
    badge = add_shape(slide, x + Inches(0.3), y + Inches(0.7), Inches(2.3), Inches(0.35),
                      RGBColor(0x0F, 0x17, 0x2A), color, Pt(1))
    add_text_box(slide, x + Inches(0.3), y + Inches(0.72), Inches(2.3), Inches(0.35),
                 mime, font_size=10, color=color, alignment=PP_ALIGN.CENTER)
    # Code example
    code_bg = add_shape(slide, x + Inches(0.15), y + Inches(1.3), Inches(2.6), Inches(3.2),
                        RGBColor(0x0D, 0x11, 0x1E))
    add_text_box(slide, x + Inches(0.25), y + Inches(1.4), Inches(2.4), Inches(3.0),
                 example, font_size=10, color=TEXT_LIGHT, font_name="Consolas")

add_footer_line(slide)
add_slide_number(slide, 12, TOTAL_SLIDES)

# ──────────────────────────────────────────────────────────────────────
# FOLIE 13: OData
# ──────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_header(slide, "OData v4 — ProductODataService", "Standardisierte Abfrage-Syntax für Produkte")

add_text_box(slide, Inches(0.75), Inches(1.4), Inches(8.5), Inches(0.7),
             "OData (Open Data Protocol) ermöglicht SQL-ähnliche Abfragen über REST. "
             "Der ProductODataService bietet $filter, $orderby, $top, $skip, $select und $count.",
             font_size=14, color=TEXT_LIGHT)

queries_data = [
    ["Abfrage", "URL"],
    ["Alle Produkte", "/odata/Products"],
    ["Preis > 500€", "/odata/Products?$filter=Price gt 500"],
    ["Top 5 teuerste", "/odata/Products?$orderby=Price desc&$top=5"],
    ["Nur Laptops", "/odata/Products?$filter=Category eq 'Laptops'"],
    ["Name + Preis", "/odata/Products?$select=Name,Price"],
    ["Mit Anzahl", "/odata/Products?$count=true&$top=3"],
]

add_table_slide(slide, Inches(0.5), Inches(2.4), Inches(9.0), queries_data,
                [Inches(2.5), Inches(6.5)], ACCENT_PURPLE)

add_text_box(slide, Inches(0.75), Inches(5.8), Inches(8.5), Inches(0.8),
             "Projektbezug most wanTED: OData kann dazu verwendet werden, "
             "die Ergebnisse abgegebener Fragebögen generisch über ein SQL-ähnliches "
             "Abfrageinterface bereitzustellen.",
             font_size=12, color=TEXT_GRAY)

add_footer_line(slide)
add_slide_number(slide, 13, TOTAL_SLIDES)

# ──────────────────────────────────────────────────────────────────────
# FOLIE 14: gRPC Logging
# ──────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_header(slide, "Zentrales Logging — gRPC", "LoggingService mit Protocol Buffers")

# Proto Definition
add_text_box(slide, Inches(0.5), Inches(1.5), Inches(4.3), Inches(0.3),
             "Protobuf-Definition (logging.proto)", font_size=14, color=ACCENT_CYAN, bold=True)

proto_code = (
    'service LoggingGrpcService {\n'
    '  rpc LogError (LogEntry)\n'
    '      returns (LogResponse);\n'
    '  rpc GetLogs (LogFilter)\n'
    '      returns (stream LogEntry);\n'
    '}\n\n'
    'message LogEntry {\n'
    '  string timestamp = 1;\n'
    '  string service_name = 2;\n'
    '  string instance_url = 3;\n'
    '  string error_message = 4;\n'
    '  string http_status_code = 5;\n'
    '  int32 retry_attempt = 6;\n'
    '  string correlation_id = 7;\n'
    '}'
)

code_bg = add_shape(slide, Inches(0.5), Inches(1.9), Inches(4.3), Inches(4.4), RGBColor(0x0D, 0x11, 0x1E))
add_text_box(slide, Inches(0.7), Inches(2.0), Inches(3.9), Inches(4.2),
             proto_code, font_size=11, color=TEXT_LIGHT, font_name="Consolas")

# Benefits
add_text_box(slide, Inches(5.2), Inches(1.5), Inches(4.3), Inches(0.3),
             "Vorteile von gRPC", font_size=14, color=ACCENT_CYAN, bold=True)

benefits = [
    ("⚡ Performant", "Binäres Protobuf-Format statt\ntextbasiertem JSON — kleiner & schneller"),
    ("📋 Typisiert", "Code-Generierung aus .proto-Dateien\n→ compile-time Sicherheit"),
    ("🔄 Streaming", "Server-Streaming für\nEchtzeit-Log-Abfragen"),
    ("🛡️ Fire & Forget", "Logging-Fehler blockieren\nnicht den Hauptprozess"),
]

for i, (title, desc) in enumerate(benefits):
    y = Inches(1.9 + i * 1.1)
    card = add_shape(slide, Inches(5.2), y, Inches(4.3), Inches(0.95), BG_CARD, ACCENT_CYAN, Pt(1))
    add_text_box(slide, Inches(5.35), y + Inches(0.08), Inches(4.0), Inches(0.3),
                 title, font_size=13, color=ACCENT_CYAN, bold=True)
    add_text_box(slide, Inches(5.35), y + Inches(0.4), Inches(4.0), Inches(0.5),
                 desc, font_size=11, color=TEXT_LIGHT)

add_footer_line(slide)
add_slide_number(slide, 14, TOTAL_SLIDES)

# ──────────────────────────────────────────────────────────────────────
# FOLIE 15: Webhooks
# ──────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_header(slide, "Webhook Pattern", "Event-basierte Benachrichtigungen bei Produktänderungen")

# Flow
flow_steps = [
    ("1. Subscribe", "WebhookSubscriber\nregistriert Callback-URL\nbeim PaymentService", ACCENT_BLUE),
    ("2. Event", "Neues Payment wird\nerstellt (POST)\n→ Event ausgelöst", ACCENT_ORANGE),
    ("3. Notify", "PaymentService sendet\nHTTP POST mit Payload\nan alle Subscriber", ACCENT_GREEN),
    ("4. Process", "WebhookSubscriber\nverarbeitet die\nBenachrichtigung", ACCENT_PURPLE),
]

for i, (step, desc, color) in enumerate(flow_steps):
    x = Inches(0.35 + i * 2.4)
    card = add_shape(slide, x, Inches(1.6), Inches(2.15), Inches(2.5), BG_CARD, color, Pt(1))
    add_accent_bar(slide, x, Inches(1.6), Inches(2.15), Pt(3), color)
    add_text_box(slide, x + Inches(0.1), Inches(1.8), Inches(1.95), Inches(0.4),
                 step, font_size=18, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), Inches(2.4), Inches(1.95), Inches(1.5),
                 desc, font_size=12, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

# Payload example
add_text_box(slide, Inches(0.5), Inches(4.4), Inches(9.0), Inches(0.3),
             "Webhook-Payload (Beispiel)", font_size=14, color=ACCENT_GREEN, bold=True)

payload = ('{ "EventType": "payment.created",\n'
           '  "Data": { "Id": 42, "Amount": 149.99, "Currency": "EUR" },\n'
           '  "Timestamp": "2026-06-08T10:30:00Z" }')

code_bg = add_shape(slide, Inches(0.5), Inches(4.8), Inches(9.0), Inches(1.4), RGBColor(0x0D, 0x11, 0x1E))
add_text_box(slide, Inches(0.7), Inches(4.9), Inches(8.6), Inches(1.2),
             payload, font_size=13, color=TEXT_LIGHT, font_name="Consolas")

add_footer_line(slide)
add_slide_number(slide, 15, TOTAL_SLIDES)

# ──────────────────────────────────────────────────────────────────────
# FOLIE 16: Fachartikelanalyse
# ──────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_header(slide, "Fachartikelanalyse", "Montesi & Weber (2016) — Circuit Breakers, Discovery, API Gateways")

# 3 Patterns
patterns_fa = [
    ("Circuit Breaker", "Verhindert kaskadierende Fehler\nzwischen Services.\n\nUmsetzung: Polly im MeiShop\nfür CreditCard-Aufrufe", ACCENT_RED),
    ("Service Discovery", "Dynamische Skalierung ohne\nmanuelle Konfigurationsänderungen.\n\nUmsetzung: HashiCorp Consul\n— alle Services registrieren sich", ACCENT_BLUE),
    ("API Gateway", "Zentraler Einstiegspunkt mit\nCross-Cutting Concerns.\n\nUmsetzung: MeiShop als\nGateway + Router", ACCENT_GREEN),
]

for i, (title, desc, color) in enumerate(patterns_fa):
    x = Inches(0.4 + i * 3.15)
    card = add_shape(slide, x, Inches(1.5), Inches(2.9), Inches(3.5), BG_CARD, color, Pt(1))
    add_accent_bar(slide, x, Inches(1.5), Inches(2.9), Pt(3), color)
    add_text_box(slide, x + Inches(0.15), Inches(1.7), Inches(2.6), Inches(0.5),
                 title, font_size=18, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.15), Inches(2.3), Inches(2.6), Inches(2.5),
                 desc, font_size=12, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

# Reflexion
add_shape(slide, Inches(0.4), Inches(5.3), Inches(9.2), Inches(1.2), BG_CARD, ACCENT_PURPLE, Pt(1))
add_text_box(slide, Inches(0.6), Inches(5.4), Inches(8.8), Inches(0.3),
             "Reflexion", font_size=14, color=ACCENT_PURPLE, bold=True)
add_text_box(slide, Inches(0.6), Inches(5.75), Inches(8.8), Inches(0.7),
             "Nutzen: Bewährte Lösungen für verteilte Systeme, gute .NET-Integration. "
             "Grenzen: Vernachlässigt organisatorische Aspekte und neuere Konzepte wie "
             "Service Mesh (Istio, Linkerd).", font_size=12, color=TEXT_LIGHT)

add_footer_line(slide)
add_slide_number(slide, 16, TOTAL_SLIDES)

# ──────────────────────────────────────────────────────────────────────
# FOLIE 17: SOAP/Workflow & BPEL (TED 8)
# ──────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_header(slide, "Produktempfehlung — SOAP & Workflow", "TED 8: Genehmigungsprozess mit Workflow-Engine")

# Workflow steps
wf_steps = [
    ("SOAP\nEndpoint", "XML-Payload\nempfangen", ACCENT_BLUE),
    ("Validierung", "Produkt prüfen\n(ProductService)", ACCENT_CYAN),
    ("Auto-Check", "Preisschwelle\nprüfen", ACCENT_GREEN),
    ("Human\nTask", "Gesellschafter\nFreigabe", ACCENT_ORANGE),
    ("Ergebnis", "Genehmigt oder\nAbgelehnt", ACCENT_PURPLE),
]

for i, (title, desc, color) in enumerate(wf_steps):
    x = Inches(0.2 + i * 1.95)
    card = add_shape(slide, x, Inches(1.5), Inches(1.75), Inches(1.5), BG_CARD, color, Pt(1))
    add_text_box(slide, x, Inches(1.55), Inches(1.75), Inches(0.7),
                 title, font_size=13, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, Inches(2.2), Inches(1.75), Inches(0.7),
                 desc, font_size=11, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)
    if i < 4:
        add_text_box(slide, x + Inches(1.75), Inches(1.9), Inches(0.2), Inches(0.4),
                     "→", font_size=14, color=color)

# Vergleichstabelle
cmp_data = [
    ["Kriterium", "BPEL (SOA)", "Orchestrierung (Saga)", "Choreographie"],
    ["Paradigma", "Zentral, SOAP/XML", "Zentral, REST/gRPC", "Dezentral, Events"],
    ["Kopplung", "Eng", "Lose", "Sehr lose"],
    ["Fehler", "XA-Transaktionen", "Compensating Tx", "Event-basiert"],
    ["Fazit", "❌ Veraltet", "✅ Empfohlen", "⚠️ Für simple Flows"],
]

add_table_slide(slide, Inches(0.3), Inches(3.5), Inches(9.4), cmp_data,
                [Inches(1.8), Inches(2.4), Inches(2.6), Inches(2.6)])

add_footer_line(slide)
add_slide_number(slide, 17, TOTAL_SLIDES)

# ──────────────────────────────────────────────────────────────────────
# FOLIE 18: KI, Low-Code & Vision (TED 9)
# ──────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_header(slide, "KI, Low-Code & Visionäre Weiterentwicklung", "TED 9: Zukunftsperspektiven")

sections = [
    ("a) KI-Unterstützung", "LLM für automatisierte Fragebogen-\nAuswertung (Sentiment-Analyse,\nKategorisierung). Risiken:\nHalluzinationen, DSGVO.", ACCENT_BLUE, "3P"),
    ("b) Low-Code / No-Code", "Power Automate / n8n für\nFragebogen-Workflow.\nVorteil: Schnelle Umsetzung.\nNachteil: Vendor Lock-in.", ACCENT_GREEN, "3P"),
    ("c) Vision: AI-Discovery", "ML-Modelle optimieren Routing\nund Auto-Scaling in Consul.\nSelf-Healing Infrastructure\nstatt manueller Konfiguration.", ACCENT_PURPLE, "3P"),
    ("Ergänzend: Agentic AI", "Agentic AI vs. RPA:\nRPA = regelbasiert\nAgentic AI = lernfähig, adaptiv\nFrameworks: OpenClaw", ACCENT_ORANGE, "3P"),
]

for i, (title, desc, color, pts) in enumerate(sections):
    col = i % 2
    row = i // 2
    x = Inches(0.4 + col * 4.8)
    y = Inches(1.5 + row * 2.7)
    card = add_shape(slide, x, y, Inches(4.5), Inches(2.3), BG_CARD, color, Pt(1))
    add_accent_bar(slide, x, y, Pt(4), Inches(2.3), color)
    # Points badge
    badge = add_shape(slide, x + Inches(3.6), y + Inches(0.1), Inches(0.7), Inches(0.35),
                      color)
    add_text_box(slide, x + Inches(3.6), y + Inches(0.12), Inches(0.7), Inches(0.3),
                 pts, font_size=10, color=TEXT_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.15), Inches(3.3), Inches(0.4),
                 title, font_size=15, color=color, bold=True)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.6), Inches(4.1), Inches(1.6),
                 desc, font_size=12, color=TEXT_LIGHT)

add_footer_line(slide)
add_slide_number(slide, 18, TOTAL_SLIDES)

# ──────────────────────────────────────────────────────────────────────
# FOLIE 19: Live-Demo
# ──────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_header(slide, "Live-Demo", "Gesamtlösung demonstrieren")

demo_steps = [
    ("1", "Consul Dashboard", "http://localhost:8500", "Alle registrierten Services\nund Health Checks zeigen", ACCENT_BLUE),
    ("2", "MeiShop Frontend", "http://localhost:7024", "Produktliste aus\nProductService + FtpService", ACCENT_CYAN),
    ("3", "OData Query", "localhost:7500/odata/Products\n?$filter=Price gt 500", "Standardisierte Abfrage\nmit $filter, $orderby", ACCENT_GREEN),
    ("4", "SAGA Bestellung", "POST localhost:7700\n/api/orders", "Order-Workflow mit\nCompensating Transactions", ACCENT_ORANGE),
    ("5", "Content Negotiation", "POST localhost:7400\nAccept: text/csv", "Payment in JSON, XML\noder CSV erstellen", ACCENT_PURPLE),
]

for i, (num, title, url, desc, color) in enumerate(demo_steps):
    y = Inches(1.4 + i * 1.05)
    # Number circle
    circle = add_shape(slide, Inches(0.5), y + Inches(0.1), Inches(0.55), Inches(0.55), color)
    add_text_box(slide, Inches(0.5), y + Inches(0.14), Inches(0.55), Inches(0.5),
                 num, font_size=18, color=TEXT_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Title
    add_text_box(slide, Inches(1.2), y + Inches(0.05), Inches(2.2), Inches(0.35),
                 title, font_size=15, color=color, bold=True)
    add_text_box(slide, Inches(1.2), y + Inches(0.4), Inches(2.2), Inches(0.5),
                 desc, font_size=11, color=TEXT_LIGHT)
    # URL
    url_bg = add_shape(slide, Inches(4.8), y + Inches(0.08), Inches(4.7), Inches(0.65),
                       RGBColor(0x0D, 0x11, 0x1E), color, Pt(1))
    add_text_box(slide, Inches(4.95), y + Inches(0.15), Inches(4.4), Inches(0.55),
                 url, font_size=10, color=ACCENT_CYAN, font_name="Consolas")

add_footer_line(slide)
add_slide_number(slide, 19, TOTAL_SLIDES)

# ──────────────────────────────────────────────────────────────────────
# FOLIE 20: Fazit & Ausblick
# ──────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

# Accent bar top
add_accent_bar(slide, Inches(0), Inches(0), Inches(10), Pt(4), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.6), Inches(8.4), Inches(0.6),
             "Fazit & Ausblick", font_size=30, color=TEXT_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Achieved
add_shape(slide, Inches(0.4), Inches(1.5), Inches(4.4), Inches(4.5), BG_CARD, ACCENT_GREEN, Pt(1))
add_accent_bar(slide, Inches(0.4), Inches(1.5), Inches(4.4), Pt(3), ACCENT_GREEN)
add_text_box(slide, Inches(0.6), Inches(1.7), Inches(4.0), Inches(0.4),
             "✅ Erreicht", font_size=18, color=ACCENT_GREEN, bold=True)

achieved = [
    "• 9 eigenständige Microservices",
    "• 5 Kommunikationsprotokolle",
    "• Consul Service Discovery",
    "• SAGA Pattern mit Compensation",
    "• Content Negotiation (JSON/XML/CSV)",
    "• Round-Robin Load Balancing",
    "• Polly Retry mit Exp. Backoff",
    "• gRPC Logging-Infrastruktur",
    "• Webhook Event-System",
    "• OData v4 Abfrage-Interface",
]

add_bullet_list(slide, Inches(0.6), Inches(2.2), Inches(4.0), Inches(3.5), achieved,
                font_size=12, color=TEXT_LIGHT, spacing=Pt(5))

# Outlook
add_shape(slide, Inches(5.2), Inches(1.5), Inches(4.4), Inches(4.5), BG_CARD, ACCENT_PURPLE, Pt(1))
add_accent_bar(slide, Inches(5.2), Inches(1.5), Inches(4.4), Pt(3), ACCENT_PURPLE)
add_text_box(slide, Inches(5.4), Inches(1.7), Inches(4.0), Inches(0.4),
             "🔮 Ausblick", font_size=18, color=ACCENT_PURPLE, bold=True)

outlook = [
    "• Containerisierung (Docker)",
    "• Service Mesh (Istio/Linkerd)",
    "• Async Messaging (RabbitMQ)",
    "• KI-gestützte Observability",
    "• AI-Driven Auto-Scaling",
    "• Secrets Management (Vault)",
    "• CI/CD Pipeline",
    "• End-to-End Tracing",
]

add_bullet_list(slide, Inches(5.4), Inches(2.2), Inches(4.0), Inches(3.5), outlook,
                font_size=12, color=TEXT_LIGHT, spacing=Pt(5))

# Thank you
add_text_box(slide, Inches(0.8), Inches(6.3), Inches(8.4), Inches(0.5),
             "Vielen Dank für Ihre Aufmerksamkeit!", font_size=18,
             color=ACCENT_CYAN, bold=True, alignment=PP_ALIGN.CENTER)

add_footer_line(slide)
add_slide_number(slide, 20, TOTAL_SLIDES)

# ═══════════════════════════════════════════════════════════════════════
#  SPEICHERN
# ═══════════════════════════════════════════════════════════════════════

output_path = os.path.join(os.path.dirname(__file__), "documentation", "most_wanTED_Praesentation.pptx")
prs.save(output_path)
print(f"\n[OK] Praesentation erfolgreich erstellt: {output_path}")
print(f"     {TOTAL_SLIDES} Folien generiert.")
